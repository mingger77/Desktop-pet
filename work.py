"""根据 JSON 数据生成 Office 文档（PPT / Word / Excel / PDF）。"""

import os
import json
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from openpyxl import Workbook
from fpdf import FPDF


# OpenAI function calling 工具定义
GENERATE_DOC_TOOL = {
    "type": "function",
    "function": {
        "name": "generate_document",
        "description": "根据结构化数据生成 Office 文档（PPT / Word / Excel / PDF）并保存到桌面",
        "parameters": {
            "type": "object",
            "properties": {
                "type": {
                    "type": "string",
                    "enum": ["pptx", "docx", "xlsx", "pdf"],
                    "description": "文档类型"
                },
                "title": {
                    "type": "string",
                    "description": "文档标题"
                },
                "content": {
                    "type": "array",
                    "description": "文档内容，每项包含 heading 和 body",
                    "items": {
                        "type": "object",
                        "properties": {
                            "heading": {"type": "string"},
                            "body": {"type": "string"}
                        },
                        "required": ["heading", "body"]
                    }
                }
            },
            "required": ["type", "title", "content"]
        }
    }
}


def handle_tool_call(tool_call, output_dir=None):
    """执行工具调用，返回结果消息。
    output_dir: 输出目录，默认使用桌面。"""
    name = tool_call.function.name
    if name == "generate_document":
        args = json.loads(tool_call.function.arguments)
        base_dir = output_dir or os.path.expanduser("~")
        output_path = os.path.join(base_dir, args.get("title", "文档"))
        try:
            result_path = generate_document(args, output_path)
            return {
                "role": "tool",
                "tool_call_id": tool_call.id,
                "content": f"文档已生成: {result_path}",
            }
        except Exception as e:
            return {
                "role": "tool",
                "tool_call_id": tool_call.id,
                "content": f"生成失败: {str(e)}",
            }
    return {
        "role": "tool",
        "tool_call_id": tool_call.id,
        "content": f"未知工具: {name}",
    }


def generate_pptx(data, output_path):
    """生成 PPTX 文件。"""
    prs = Presentation()
    title = data.get("title", "文档")
    content = data.get("content", [])

    if not content:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.placeholders[0].text = title
    else:
        for item in content:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = item.get("heading", "")
            slide.placeholders[1].text = item.get("body", "")

    prs.save(output_path)
    return output_path


def generate_docx(data, output_path):
    """生成 DOCX 文件。"""
    doc = Document()
    doc.add_heading(data.get("title", "文档"), 0)
    content = data.get("content", [])

    for item in content:
        heading = item.get("heading", "")
        body = item.get("body", "")
        if heading:
            doc.add_heading(heading, 1)
        if body:
            for paragraph in body.split("\n"):
                doc.add_paragraph(paragraph)

    doc.save(output_path)
    return output_path


def generate_xlsx(data, output_path):
    """生成 XLSX 文件。"""
    wb = Workbook()
    ws = wb.active
    ws.title = data.get("title", "Sheet1")

    content = data.get("content", [])
    for row_idx, item in enumerate(content, 1):
        ws.cell(row=row_idx, column=1, value=item.get("heading", ""))
        ws.cell(row=row_idx, column=2, value=item.get("body", ""))

    wb.save(output_path)
    return output_path


def generate_pdf(data, output_path):
    """生成 PDF 文件（支持中文）。"""
    pdf = FPDF()
    pdf.add_page()

    # 使用系统黑体支持中文
    font_path = "C:/Windows/Fonts/simhei.ttf"
    if os.path.exists(font_path):
        pdf.add_font("CJK", "", font_path, uni=True)
        font_name = "CJK"
    else:
        font_name = "Helvetica"

    pdf.set_font(font_name, size=16)
    pdf.cell(0, 10, text=data.get("title", "文档"), new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.ln(10)

    content = data.get("content", [])
    for item in content:
        heading = item.get("heading", "")
        body = item.get("body", "")
        if heading:
            pdf.set_font(font_name, size=13)
            pdf.cell(0, 10, text=heading, new_x="LMARGIN", new_y="NEXT")
        if body:
            pdf.set_font(font_name, size=11)
            for line in body.split("\n"):
                pdf.multi_cell(0, 8, text=line)
                pdf.ln(2)

    pdf.output(output_path)
    return output_path


def generate_document(json_data, output_path):
    """根据 JSON 数据生成文档，返回文件路径。

    JSON 格式:
    {
        "type": "pptx" | "docx" | "xlsx" | "pdf",
        "title": "文档标题",
        "content": [
            {"heading": "章节标题", "body": "正文内容"}
        ]
    }
    """
    doc_type = json_data.get("type", "pptx").lower()

    # 确保输出路径扩展名正确
    ext_map = {"pptx": ".pptx", "docx": ".docx", "xlsx": ".xlsx", "pdf": ".pdf"}
    expected_ext = ext_map.get(doc_type, f".{doc_type}")
    if not output_path.lower().endswith(expected_ext):
        output_path += expected_ext

    generators = {
        "pptx": generate_pptx,
        "docx": generate_docx,
        "xlsx": generate_xlsx,
        "pdf": generate_pdf,
    }

    gen = generators.get(doc_type)
    if not gen:
        raise ValueError(f"不支持的文档类型: {doc_type}，仅支持 pptx/docx/xlsx/pdf")

    return gen(json_data, output_path)
